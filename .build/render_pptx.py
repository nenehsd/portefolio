# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from deck_core import W, H, POP, ALFA
import textutil as TU

EMU_PT = 12700


def P(v):
    return Emu(int(round(v * EMU_PT)))


def RGB(h):
    return RGBColor.from_string(h)


def _noline(shp):
    shp.line.fill.background()


def _nofill(shp):
    shp.fill.background()


def add_textbox(sl, o):
    s = o["s"].upper() if o.get("caps") else o["s"]
    font, bold, size = o["font"], o["bold"], o["size"]
    lines = TU.wrap(s, font, bold, size, o["w"], o.get("spacing", 0.0))
    lead = size * o["leading"]
    total = lead * len(lines)
    y0 = o["y"]
    if o.get("h") and o.get("valign") == "m":
        y0 = o["y"] + (o["h"] - total) / 2.0
    elif o.get("h") and o.get("valign") == "b":
        y0 = o["y"] + o["h"] - total
    # marge interne 0 -> position fidele; on remonte d'un poil pour l'ascender
    top = y0 - size * 0.16
    tb = sl.shapes.add_textbox(P(o["x"] - 2), P(top), P(o["w"] + 4), P(total + size * 0.4))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[o["align"]]
        p.line_spacing = Pt(lead)
        r = p.add_run()
        r.text = ln
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = o.get("italic", False)
        f.color.rgb = RGB(o["color"])
        if o.get("spacing"):
            r.font._rPr.set("spc", str(int(o["spacing"] * 100)))
    return total


def render(slides, path):
    prs = Presentation()
    prs.slide_width = P(W)
    prs.slide_height = P(H)
    blank = prs.slide_layouts[6]
    for sd in slides:
        sl = prs.slides.add_slide(blank)
        for kind, o in sd.ops:
            if kind == "rect":
                shape = MSO_SHAPE.ROUNDED_RECTANGLE if o["radius"] else MSO_SHAPE.RECTANGLE
                shp = sl.shapes.add_shape(shape, P(o["x"]), P(o["y"]), P(o["w"]), P(o["h"]))
                if o["radius"]:
                    adj = min(0.5, o["radius"] / float(min(o["w"], o["h"])))
                    shp.adjustments[0] = adj
                if o["fill"]:
                    shp.fill.solid()
                    shp.fill.fore_color.rgb = RGB(o["fill"])
                    if o["alpha"] < 1:
                        from pptx.oxml.ns import qn
                        sf = shp.fill.fore_color._xFill.find(qn('a:srgbClr'))
                        a = sf.makeelement(qn('a:alpha'), {'val': str(int(o["alpha"] * 100000))})
                        sf.append(a)
                else:
                    _nofill(shp)
                if o["line"]:
                    shp.line.color.rgb = RGB(o["line"])
                    shp.line.width = Pt(o["lw"])
                else:
                    _noline(shp)
                shp.shadow.inherit = False
                shp.text_frame.text = ""
            elif kind == "ellipse":
                shp = sl.shapes.add_shape(MSO_SHAPE.OVAL, P(o["x"]), P(o["y"]), P(o["w"]), P(o["h"]))
                if o["fill"]:
                    shp.fill.solid()
                    shp.fill.fore_color.rgb = RGB(o["fill"])
                else:
                    _nofill(shp)
                if o["line"]:
                    shp.line.color.rgb = RGB(o["line"])
                    shp.line.width = Pt(o["lw"])
                else:
                    _noline(shp)
                shp.shadow.inherit = False
            elif kind == "poly":
                shp = _freeform(sl, o["pts"])
                if o["fill"]:
                    shp.fill.solid()
                    shp.fill.fore_color.rgb = RGB(o["fill"])
                else:
                    _nofill(shp)
                if o["line"]:
                    shp.line.color.rgb = RGB(o["line"])
                    shp.line.width = Pt(o["lw"])
                else:
                    _noline(shp)
                shp.shadow.inherit = False
            elif kind == "line":
                from pptx.enum.shapes import MSO_CONNECTOR
                cn = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, P(o["x1"]), P(o["y1"]),
                                             P(o["x2"]), P(o["y2"]))
                cn.line.color.rgb = RGB(o["color"])
                cn.line.width = Pt(o["lw"])
                if o["dash"]:
                    from pptx.enum.dml import MSO_LINE_DASH_STYLE
                    cn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            elif kind == "image":
                sl.shapes.add_picture(o["path"], P(o["x"]), P(o["y"]), P(o["w"]), P(o["h"]))
            elif kind == "text":
                add_textbox(sl, o)
            elif kind == "bullets":
                y = o["y"]
                for it in o["items"]:
                    if o["marker"] == "dot":
                        d = sl.shapes.add_shape(MSO_SHAPE.OVAL, P(o["x"] + 0.4),
                                                P(y + o["size"] * 0.45 - 3), P(6), P(6))
                        d.fill.solid(); d.fill.fore_color.rgb = RGB(o["mcolor"])
                        _noline(d); d.shadow.inherit = False
                        tx, tw = o["x"] + 13, o["w"] - 13
                    elif o["marker"] == "check":
                        d = sl.shapes.add_shape(MSO_SHAPE.OVAL, P(o["x"] - 0.2),
                                                P(y + o["size"] * 0.45 - 5.2), P(10.4), P(10.4))
                        d.fill.solid(); d.fill.fore_color.rgb = RGB(o["mcolor"])
                        _noline(d); d.shadow.inherit = False
                        tfd = d.text_frame
                        tfd.margin_left = tfd.margin_right = 0
                        tfd.margin_top = tfd.margin_bottom = 0
                        pr = tfd.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
                        rr = pr.add_run(); rr.text = u"\u2713"
                        rr.font.size = Pt(7); rr.font.bold = True
                        rr.font.color.rgb = RGB("FFFFFF"); rr.font.name = POP
                        tx, tw = o["x"] + 16, o["w"] - 16
                    else:
                        tx, tw = o["x"], o["w"]
                    used = add_textbox(sl, dict(
                        s=it, x=tx, y=y, w=tw, size=o["size"], font=o["font"], bold=o["bold"],
                        italic=False, color=o["color"], align="l", h=None, valign="t",
                        leading=o["leading"], spacing=0.0, caps=False))
                    y += used + o["gap"]
    prs.save(path)


def _freeform(sl, pts):
    b = sl.shapes.build_freeform(P(pts[0][0]), P(pts[0][1]))
    b.add_line_segments([(P(x), P(y)) for x, y in pts[1:]], close=True)
    return b.convert_to_shape()
